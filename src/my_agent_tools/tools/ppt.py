from __future__ import annotations

import re
from dataclasses import dataclass
from math import floor
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from my_agent_tools.specs import BlockSpec, ChartSpec, DeckSpec, MetricSpec, SlideSpec
from my_agent_tools.template_config import LayoutBinding, load_template_config
from my_agent_tools.themes import load_theme
from my_agent_tools.tools.base import ToolResult, ensure_parent_dir


BASE_SLIDE_WIDTH = 13.333
BASE_SLIDE_HEIGHT = 7.5
SLIDE_SIZES = {
    "16:9": (13.333, 7.5),
    "4:3": (10.0, 7.5),
}
CONTINUATION_SUFFIX = "（续）"
CHART_TYPE_MAP = {
    "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar_clustered": XL_CHART_TYPE.BAR_CLUSTERED,
    "line_marker": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
}


@dataclass(frozen=True)
class TextRegion:
    width_in: float
    height_in: float
    font_size_pt: int

    def chars_per_line(self) -> int:
        width_pts = self.width_in * 72
        average_char_width = self.font_size_pt * 0.56
        return max(18, floor(width_pts / max(1.0, average_char_width)))

    def line_capacity(self) -> int:
        height_pts = self.height_in * 72
        line_height = self.font_size_pt * 1.45
        return max(4, floor(height_pts / max(10.0, line_height)))

    def chars_capacity(self) -> int:
        return self.chars_per_line() * self.line_capacity()


class PptRenderer:
    def __init__(
        self,
        spec: DeckSpec,
        template_path: str | None = None,
        theme_path: str | None = None,
        template_config_path: str | None = None,
    ) -> None:
        self.spec = spec
        self.theme = load_theme(spec.theme.name, theme_path or spec.theme.path)
        source_template = template_path or spec.constraints.template_path
        self.presentation = Presentation(source_template) if source_template else Presentation()
        self.template_config = load_template_config(template_config_path or spec.constraints.template_config_path)
        self.warnings: list[str] = []
        self.slide_width_in = BASE_SLIDE_WIDTH
        self.slide_height_in = BASE_SLIDE_HEIGHT
        self._configure_presentation()

    def render(self, output_path: str | Path) -> ToolResult:
        expanded_slides = self._expand_slides(self.spec.slides)
        if self.spec.constraints.max_slides is not None and len(expanded_slides) > self.spec.constraints.max_slides:
            self.warnings.append(
                f"Expanded slide count {len(expanded_slides)} exceeds constraints.max_slides={self.spec.constraints.max_slides}"
            )

        total = len(expanded_slides)
        for index, slide_spec in enumerate(expanded_slides, start=1):
            self._render_slide(slide_spec, slide_number=index, total_slides=total)

        output = ensure_parent_dir(output_path)
        self.presentation.save(output)
        return ToolResult(
            ok=True,
            message=f"Rendered {total} slides to {output}",
            output_path=str(output),
            warnings=self.warnings,
            meta={
                "slides_requested": len(self.spec.slides),
                "slides_rendered": total,
                "theme": self.theme.name,
                "ratio": self.spec.meta.ratio,
                "template_mode": bool(self.template_config),
            },
        )

    def _configure_presentation(self) -> None:
        width_in, height_in = SLIDE_SIZES[self.spec.meta.ratio]
        self.slide_width_in = width_in
        self.slide_height_in = height_in
        self.presentation.slide_width = Inches(width_in)
        self.presentation.slide_height = Inches(height_in)
        core = self.presentation.core_properties
        core.title = self.spec.meta.title
        if self.spec.meta.author:
            core.author = self.spec.meta.author
        if self.spec.meta.company:
            core.company = self.spec.meta.company
        if self.spec.meta.language:
            core.language = self.spec.meta.language

    def _expand_slides(self, slides: list[SlideSpec]) -> list[SlideSpec]:
        expanded: list[SlideSpec] = []
        for slide in slides:
            if slide.kind in {"content", "closing"}:
                expanded.extend(self._expand_textual_slide(slide))
            elif slide.kind == "table":
                expanded.extend(self._expand_table_slide(slide))
            else:
                expanded.append(slide)
        return expanded

    def _expand_textual_slide(self, slide: SlideSpec) -> list[SlideSpec]:
        region = self._text_region_for_slide(slide)
        limit_lines = min(region.line_capacity(), self.spec.constraints.max_text_line_units_per_slide)
        limit_chars = min(region.chars_capacity(), self.spec.constraints.max_chars_per_text_slide)
        units = self._explode_text_blocks(slide.blocks, chunk_limit=max(140, min(260, region.chars_per_line() * 4)))
        if not units:
            return [slide]

        pages: list[list[BlockSpec]] = []
        current: list[BlockSpec] = []
        current_lines = 0
        current_chars = 0
        for unit in units:
            unit_lines = self._estimated_line_units(unit, region)
            unit_chars = self._block_char_count(unit)
            if current and (current_lines + unit_lines > limit_lines or current_chars + unit_chars > limit_chars):
                pages.append(self._collapse_text_units(current))
                current = []
                current_lines = 0
                current_chars = 0
            current.append(unit)
            current_lines += unit_lines
            current_chars += unit_chars

        if current:
            pages.append(self._collapse_text_units(current))

        if len(pages) == 1:
            return [slide]

        self.warnings.append(f"Slide '{slide.title}' overflowed and was split into {len(pages)} slides")
        return [
            self._continuation_slide(slide, page_blocks, page_index=index)
            for index, page_blocks in enumerate(pages)
        ]

    def _expand_table_slide(self, slide: SlideSpec) -> list[SlideSpec]:
        table_block = next((block for block in slide.blocks if block.type == "table"), None)
        if table_block is None:
            return [slide]

        region = self._table_region_for_slide(slide)
        rows = table_block.rows or []
        estimated_rows = max(3, floor((region.height_in * 72 - 26) / 24))
        per_slide = min(self.spec.constraints.max_table_rows_per_slide, estimated_rows)
        if len(rows) <= per_slide:
            return [slide]

        pages: list[SlideSpec] = []
        for index, start in enumerate(range(0, len(rows), per_slide)):
            page_rows = rows[start : start + per_slide]
            page_blocks = [
                table_block.model_copy(
                    update={
                        "rows": page_rows,
                    }
                )
            ]
            pages.append(self._continuation_slide(slide, page_blocks, page_index=index))

        self.warnings.append(f"Table slide '{slide.title}' overflowed and was split into {len(pages)} slides")
        return pages

    def _explode_text_blocks(self, blocks: list[BlockSpec], chunk_limit: int) -> list[BlockSpec]:
        exploded: list[BlockSpec] = []
        for block in blocks:
            if block.type == "bullet_list":
                for item in block.items or []:
                    exploded.append(BlockSpec(type="bullet_list", items=[item]))
            elif block.type == "paragraph":
                for chunk in self._split_text_chunks(block.text or "", limit=chunk_limit):
                    exploded.append(BlockSpec(type="paragraph", text=chunk))
            else:
                exploded.append(block)
        return exploded

    def _collapse_text_units(self, units: list[BlockSpec]) -> list[BlockSpec]:
        collapsed: list[BlockSpec] = []
        current_bullets: list[str] = []
        for unit in units:
            if unit.type == "bullet_list":
                current_bullets.extend(unit.items or [])
                continue

            if current_bullets:
                collapsed.append(BlockSpec(type="bullet_list", items=current_bullets.copy()))
                current_bullets = []
            collapsed.append(unit)

        if current_bullets:
            collapsed.append(BlockSpec(type="bullet_list", items=current_bullets))

        return collapsed

    def _continuation_slide(self, slide: SlideSpec, blocks: list[BlockSpec], page_index: int) -> SlideSpec:
        is_continuation = page_index > 0
        title = slide.title if not is_continuation else f"{slide.title}{CONTINUATION_SUFFIX}"
        notes = slide.speaker_notes
        if is_continuation and notes:
            notes = f"{notes}\n\n[Continuation page {page_index + 1}]"
        return slide.model_copy(
            update={
                "title": title,
                "subtitle": slide.subtitle if page_index == 0 else None,
                "blocks": blocks,
                "speaker_notes": notes,
                "metadata": {
                    **slide.metadata,
                    "is_continuation": is_continuation,
                    "continuation_index": page_index,
                },
            }
        )

    def _split_text_chunks(self, text: str, limit: int) -> list[str]:
        compact = " ".join(text.split())
        if len(compact) <= limit:
            return [compact] if compact else []

        pieces = [piece.strip() for piece in re.split(r"(?<=[。！？!?\.])\s+", compact) if piece.strip()]
        if len(pieces) == 1:
            return [compact[i : i + limit] for i in range(0, len(compact), limit)]

        chunks: list[str] = []
        current = ""
        for piece in pieces:
            candidate = piece if not current else f"{current} {piece}"
            if current and len(candidate) > limit:
                chunks.append(current)
                current = piece
            else:
                current = candidate
        if current:
            chunks.append(current)
        return chunks

    def _block_char_count(self, block: BlockSpec) -> int:
        if block.type == "paragraph":
            return len(block.text or "")
        if block.type == "bullet_list":
            return sum(len(item) for item in (block.items or []))
        if block.type == "code":
            return len(block.content or "")
        if block.type == "table":
            return sum(len(str(cell)) for row in (block.rows or []) for cell in row)
        return 120

    def _estimated_line_units(self, block: BlockSpec, region: TextRegion) -> int:
        chars_per_line = region.chars_per_line()
        if block.type == "paragraph":
            return max(1, floor(len(block.text or "") / chars_per_line) + 1) + 1
        if block.type == "bullet_list":
            return sum(max(1, floor(len(item) / max(12, chars_per_line - 6)) + 1) for item in (block.items or []))
        if block.type == "code":
            return max(3, floor(len(block.content or "") / max(24, chars_per_line - 10)) + 2)
        if block.type == "table":
            return len(block.rows or []) + 2
        return 4

    def _render_slide(self, slide_spec: SlideSpec, slide_number: int, total_slides: int) -> None:
        binding = self._binding_for(slide_spec.kind)
        layout = self._resolve_layout(binding) if binding else self._blank_layout()
        slide = self.presentation.slides.add_slide(layout)
        if binding:
            self._render_template_slide(slide, slide_spec, binding)
        else:
            self._paint_slide_background(slide)
            self._render_manual_slide(slide, slide_spec)

        self._render_footer(slide, slide_number, total_slides, binding)
        self._apply_speaker_notes(slide, slide_spec)

    def _blank_layout(self):
        for layout in reversed(self.presentation.slide_layouts):
            if len(layout.placeholders) == 0:
                return layout
        return self.presentation.slide_layouts[-1]

    def _binding_for(self, slide_kind: str) -> LayoutBinding | None:
        if not self.template_config:
            return None
        return self.template_config.binding_for(slide_kind)

    def _resolve_layout(self, binding: LayoutBinding):
        if binding.layout_index is not None:
            return self.presentation.slide_layouts[binding.layout_index]
        for layout in self.presentation.slide_layouts:
            if layout.name == binding.layout_name:
                return layout
        self.warnings.append(f"Template layout '{binding.layout_name}' not found, using blank layout")
        return self._blank_layout()

    def _placeholder_by_idx(self, slide, placeholder_idx: int | None):
        if placeholder_idx is None:
            return None
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == placeholder_idx:
                return placeholder
        return None

    def _layout_placeholder(self, slide_kind: str, placeholder_idx: int | None):
        binding = self._binding_for(slide_kind)
        if not binding or placeholder_idx is None:
            return None
        layout = self._resolve_layout(binding)
        for placeholder in layout.placeholders:
            if placeholder.placeholder_format.idx == placeholder_idx:
                return placeholder
        return None

    def _shape_region(self, slide, placeholder_idx: int):
        shape = self._placeholder_by_idx(slide, placeholder_idx)
        if shape is None:
            return None
        return (
            shape.left / 914400,
            shape.top / 914400,
            shape.width / 914400,
            shape.height / 914400,
        )

    def _text_region_for_slide(self, slide: SlideSpec) -> TextRegion:
        binding = self._binding_for(slide.kind)
        if binding and binding.body_placeholder_idxs:
            placeholder = self._layout_placeholder(slide.kind, binding.body_placeholder_idxs[0])
            if placeholder is not None:
                return TextRegion(
                    width_in=placeholder.width / 914400,
                    height_in=placeholder.height / 914400,
                    font_size_pt=self.theme.body_font_size,
                )
        if slide.kind == "closing":
            return TextRegion(10.9, 3.95, self.theme.body_font_size)
        return TextRegion(11.95, 5.15, self.theme.body_font_size)

    def _table_region_for_slide(self, slide: SlideSpec) -> TextRegion:
        binding = self._binding_for(slide.kind)
        if binding and binding.body_placeholder_idxs:
            placeholder = self._layout_placeholder(slide.kind, binding.body_placeholder_idxs[0])
            if placeholder is not None:
                return TextRegion(
                    width_in=placeholder.width / 914400,
                    height_in=placeholder.height / 914400,
                    font_size_pt=12,
                )
        return TextRegion(11.9, 4.85, 12)

    def _render_manual_slide(self, slide, slide_spec: SlideSpec) -> None:
        self._render_header(slide, slide_spec)
        if slide_spec.kind == "title":
            self._render_title_slide(slide, slide_spec)
        elif slide_spec.kind == "section":
            self._render_section_slide(slide, slide_spec)
        elif slide_spec.kind == "content":
            self._render_content_slide(slide, slide_spec)
        elif slide_spec.kind == "two_column":
            self._render_two_column_slide(slide, slide_spec)
        elif slide_spec.kind == "table":
            self._render_table_slide(slide, slide_spec)
        elif slide_spec.kind == "image":
            self._render_image_slide(slide, slide_spec)
        elif slide_spec.kind == "metrics":
            self._render_metrics_slide(slide, slide_spec)
        elif slide_spec.kind == "chart":
            self._render_chart_slide(slide, slide_spec)
        elif slide_spec.kind == "closing":
            self._render_closing_slide(slide, slide_spec)
        else:
            self.warnings.append(f"Unsupported slide kind '{slide_spec.kind}' on '{slide_spec.title}'")

    def _paint_slide_background(self, slide) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self._rgb(self.theme.background_color)

    def _render_header(self, slide, slide_spec: SlideSpec) -> None:
        if slide_spec.kind in {"title", "section"}:
            return
        title_box = slide.shapes.add_textbox(self._x(0.7), self._y(0.45), self._w(12.0), self._h(0.8))
        frame = title_box.text_frame
        frame.clear()
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        paragraph.text = slide_spec.title
        paragraph.alignment = PP_ALIGN.LEFT
        if paragraph.runs:
            self._apply_run_style(paragraph.runs[0], self.theme.font_family, self.theme.title_font_size, self.theme.title_color, bold=True)

        if slide_spec.subtitle:
            subtitle_box = slide.shapes.add_textbox(self._x(0.72), self._y(1.10), self._w(11.6), self._h(0.35))
            frame = subtitle_box.text_frame
            frame.clear()
            paragraph = frame.paragraphs[0]
            paragraph.text = slide_spec.subtitle
            if paragraph.runs:
                self._apply_run_style(paragraph.runs[0], self.theme.font_family, self.theme.subtitle_font_size, self.theme.muted_text_color)

    def _render_template_slide(self, slide, slide_spec: SlideSpec, binding: LayoutBinding) -> None:
        self._fill_text_placeholder(slide, binding.title_placeholder_idx, slide_spec.title, size=self.theme.title_font_size, bold=True)
        if slide_spec.subtitle and binding.subtitle_placeholder_idx is not None:
            self._fill_text_placeholder(slide, binding.subtitle_placeholder_idx, slide_spec.subtitle, size=self.theme.subtitle_font_size)

        if slide_spec.kind in {"title", "section"}:
            if slide_spec.kind == "section" and binding.body_placeholder_idxs:
                self._fill_text_placeholder(
                    slide,
                    binding.body_placeholder_idxs[0],
                    slide_spec.subtitle or "Section",
                    size=self.theme.subtitle_font_size,
                )
            return

        if slide_spec.kind in {"content", "closing"} and binding.body_placeholder_idxs:
            body = self._placeholder_by_idx(slide, binding.body_placeholder_idxs[0])
            if body is not None:
                self._render_blocks_into_text_frame(body.text_frame, slide_spec.blocks)
                return

        if slide_spec.kind == "two_column" and len(binding.body_placeholder_idxs) >= 2:
            left = self._placeholder_by_idx(slide, binding.body_placeholder_idxs[0])
            right = self._placeholder_by_idx(slide, binding.body_placeholder_idxs[1])
            if left is not None and right is not None:
                self._render_blocks_into_text_frame(left.text_frame, [slide_spec.blocks[0]])
                self._render_blocks_into_text_frame(right.text_frame, [slide_spec.blocks[1]])
                return

        if slide_spec.kind == "image" and binding.picture_placeholder_idx is not None:
            image_block = next(block for block in slide_spec.blocks if block.type == "image")
            picture_placeholder = self._placeholder_by_idx(slide, binding.picture_placeholder_idx)
            if picture_placeholder is not None and Path(image_block.path or "").exists():
                try:
                    picture_placeholder.insert_picture(str(image_block.path))
                    return
                except Exception:
                    self.warnings.append(f"Template picture placeholder insert failed on '{slide_spec.title}', fallback to manual image render")

        if slide_spec.kind == "table":
            region = self._shape_region(slide, binding.body_placeholder_idxs[0]) if binding.body_placeholder_idxs else None
            self._render_table_slide(slide, slide_spec, region=region)
            return

        if slide_spec.kind == "chart":
            region = self._shape_region(slide, binding.body_placeholder_idxs[0]) if binding.body_placeholder_idxs else None
            self._render_chart_slide(slide, slide_spec, region=region)
            return

        if slide_spec.kind == "metrics":
            region = self._shape_region(slide, binding.body_placeholder_idxs[0]) if binding.body_placeholder_idxs else None
            self._render_metrics_slide(slide, slide_spec, region=region)
            return

        self._render_manual_slide(slide, slide_spec)

    def _render_footer(self, slide, slide_number: int, total_slides: int, binding: LayoutBinding | None) -> None:
        footer_text = " | ".join(part for part in [self.spec.meta.company, self.spec.meta.author] if part)
        page_text = f"{slide_number} / {total_slides}"
        if binding and binding.use_template_footer:
            footer_placeholder = self._placeholder_by_idx(slide, binding.footer_placeholder_idx)
            page_placeholder = self._placeholder_by_idx(slide, binding.slide_number_placeholder_idx)
            rendered = False
            if footer_placeholder is not None:
                footer_placeholder.text = footer_text
                rendered = True
            if page_placeholder is not None:
                page_placeholder.text = page_text
                rendered = True
            if rendered:
                return

        page_box = slide.shapes.add_textbox(self._x(0.7), self._y(7.02), self._w(1.2), self._h(0.2))
        page_frame = page_box.text_frame
        page_frame.clear()
        page_para = page_frame.paragraphs[0]
        page_para.text = page_text
        if page_para.runs:
            self._apply_run_style(page_para.runs[0], self.theme.font_family, self.theme.small_font_size, self.theme.muted_text_color, bold=True)

        footer_box = slide.shapes.add_textbox(self._x(2.0), self._y(7.02), self._w(10.7), self._h(0.2))
        frame = footer_box.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = footer_text
        paragraph.alignment = PP_ALIGN.RIGHT
        if paragraph.runs:
            self._apply_run_style(paragraph.runs[0], self.theme.font_family, self.theme.small_font_size, self.theme.muted_text_color)

    def _fill_text_placeholder(
        self,
        slide,
        placeholder_idx: int | None,
        text: str,
        *,
        size: int,
        bold: bool = False,
    ) -> bool:
        placeholder = self._placeholder_by_idx(slide, placeholder_idx)
        if placeholder is None:
            return False
        frame = placeholder.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        if paragraph.runs:
            self._apply_run_style(
                paragraph.runs[0],
                self.theme.font_family,
                size,
                self.theme.title_color if bold else self.theme.text_color,
                bold=bold,
            )
        return True

    def _render_title_slide(self, slide, slide_spec: SlideSpec) -> None:
        accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, self._x(0), self._y(0), self._w(0.24), self._h(7.5))
        accent.fill.solid()
        accent.fill.fore_color.rgb = self._rgb(self.theme.accent_color)
        accent.line.fill.background()

        title_box = slide.shapes.add_textbox(self._x(1.1), self._y(1.55), self._w(10.5), self._h(1.5))
        frame = title_box.text_frame
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.text = slide_spec.title
        if paragraph.runs:
            self._apply_run_style(paragraph.runs[0], self.theme.font_family, 30, self.theme.title_color, bold=True)

        subtitle = slide_spec.subtitle or self.spec.meta.subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(self._x(1.15), self._y(3.05), self._w(9.5), self._h(0.6))
            frame = subtitle_box.text_frame
            frame.clear()
            paragraph = frame.paragraphs[0]
            paragraph.text = subtitle
            if paragraph.runs:
                self._apply_run_style(paragraph.runs[0], self.theme.font_family, 17, self.theme.muted_text_color)

    def _render_section_slide(self, slide, slide_spec: SlideSpec) -> None:
        band = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            self._x(0.9),
            self._y(1.9),
            self._w(11.2),
            self._h(2.6),
        )
        band.fill.solid()
        band.fill.fore_color.rgb = self._rgb(self.theme.accent_soft_color)
        band.line.fill.background()

        title_box = slide.shapes.add_textbox(self._x(1.25), self._y(2.55), self._w(9.5), self._h(0.8))
        frame = title_box.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = slide_spec.title
        if paragraph.runs:
            self._apply_run_style(paragraph.runs[0], self.theme.font_family, 26, self.theme.title_color, bold=True)

        subtitle = slide_spec.subtitle or "Section"
        subtitle_box = slide.shapes.add_textbox(self._x(1.25), self._y(3.3), self._w(6.0), self._h(0.4))
        frame = subtitle_box.text_frame
        paragraph = frame.paragraphs[0]
        paragraph.text = subtitle
        if paragraph.runs:
            self._apply_run_style(paragraph.runs[0], self.theme.font_family, self.theme.subtitle_font_size, self.theme.muted_text_color)

    def _render_content_slide(self, slide, slide_spec: SlideSpec) -> None:
        content_box = self._add_content_panel(slide, 0.7, 1.55, 11.95, 5.15)
        self._render_blocks_into_text_frame(content_box.text_frame, slide_spec.blocks)

    def _render_two_column_slide(self, slide, slide_spec: SlideSpec) -> None:
        left = self._add_content_panel(slide, 0.7, 1.55, 5.7, 5.15)
        right = self._add_content_panel(slide, 6.95, 1.55, 5.7, 5.15)
        self._render_blocks_into_text_frame(left.text_frame, [slide_spec.blocks[0]])
        self._render_blocks_into_text_frame(right.text_frame, [slide_spec.blocks[1]])

        for extra in slide_spec.blocks[2:]:
            self.warnings.append(f"Extra block on two_column slide '{slide_spec.title}' was ignored: {extra.type}")

    def _render_table_slide(self, slide, slide_spec: SlideSpec, region: tuple[float, float, float, float] | None = None) -> None:
        table_block = next(block for block in slide_spec.blocks if block.type == "table")
        rows = len(table_block.rows or []) + 1
        cols = len(table_block.columns or [])
        left, top, width, height = region if region is not None else (0.7, 1.7, 11.9, 4.85)
        shape = slide.shapes.add_table(rows, cols, self._x(left), self._y(top), self._w(width), self._h(height))
        table = shape.table

        for col_idx, heading in enumerate(table_block.columns or []):
            cell = table.cell(0, col_idx)
            cell.text = str(heading)
            self._format_table_cell(cell, header=True)

        for row_idx, row in enumerate(table_block.rows or [], start=1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)
                self._format_table_cell(cell, header=False, stripe=(row_idx % 2 == 0))

    def _render_image_slide(self, slide, slide_spec: SlideSpec) -> None:
        image_block = next(block for block in slide_spec.blocks if block.type == "image")
        image_path = Path(image_block.path or "")
        if not image_path.exists():
            self.warnings.append(f"Image not found for slide '{slide_spec.title}': {image_path}")
            placeholder = self._add_content_panel(slide, 0.7, 1.7, 11.95, 4.9)
            self._render_blocks_into_text_frame(
                placeholder.text_frame,
                [BlockSpec(type="paragraph", text=f"Missing image: {image_path}")],
            )
            return
        slide.shapes.add_picture(str(image_path), self._x(0.95), self._y(1.7), width=self._w(11.4), height=self._h(4.9))

    def _render_metrics_slide(self, slide, slide_spec: SlideSpec, region: tuple[float, float, float, float] | None = None) -> None:
        block = next(block for block in slide_spec.blocks if block.type == "metrics")
        metrics = block.metrics or []
        if not metrics:
            self.warnings.append(f"Metrics slide '{slide_spec.title}' had no metrics payload")
            return

        left, top, total_width, height = region if region is not None else (0.7, 2.0, 11.95, 3.1)
        card_gap = 0.25
        card_width = (total_width - card_gap * (len(metrics) - 1)) / max(1, len(metrics))
        current_left = left
        for metric in metrics:
            self._render_metric_card(slide, metric, current_left, top, card_width, height)
            current_left += card_width + card_gap

    def _render_chart_slide(self, slide, slide_spec: SlideSpec, region: tuple[float, float, float, float] | None = None) -> None:
        chart_block = next(block for block in slide_spec.blocks if block.type == "chart")
        if region is None:
            self._add_content_panel(slide, 0.7, 1.55, 11.95, 5.15)
            left, top, width, height = 1.0, 1.9, 11.3, 4.35
        else:
            left, top, width, height = region
        chart_data = self._build_chart_data(chart_block.chart)
        chart_shape = slide.shapes.add_chart(
            CHART_TYPE_MAP[chart_block.chart.chart_type],
            self._x(left),
            self._y(top),
            self._w(width),
            self._h(height),
            chart_data,
        )
        chart = chart_shape.chart
        chart.has_legend = len(chart_block.chart.series) > 1
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        if hasattr(chart, "value_axis"):
            chart.value_axis.has_major_gridlines = True
            chart.value_axis.tick_labels.font.size = Pt(10)
        if hasattr(chart, "category_axis"):
            chart.category_axis.tick_labels.font.size = Pt(11)
        if chart_block.chart.title:
            chart.has_title = True
            chart.chart_title.text_frame.text = chart_block.chart.title

    def _build_chart_data(self, chart_spec: ChartSpec) -> CategoryChartData:
        data = CategoryChartData()
        data.categories = chart_spec.categories
        for series in chart_spec.series:
            data.add_series(series.name, series.values)
        return data

    def _render_closing_slide(self, slide, slide_spec: SlideSpec) -> None:
        content_box = self._add_content_panel(slide, 1.1, 1.75, 10.9, 3.95)
        self._render_blocks_into_text_frame(content_box.text_frame, slide_spec.blocks)
        thanks_box = slide.shapes.add_textbox(self._x(1.12), self._y(5.95), self._w(4.0), self._h(0.5))
        frame = thanks_box.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = "Thank you"
        if paragraph.runs:
            self._apply_run_style(paragraph.runs[0], self.theme.font_family, 22, self.theme.accent_color, bold=True)

    def _render_blocks_into_text_frame(self, text_frame, blocks: Iterable[BlockSpec]) -> None:
        text_frame.clear()
        text_frame.word_wrap = True
        first_paragraph_used = False

        for block in blocks:
            if block.type == "paragraph":
                paragraph = text_frame.paragraphs[0] if not first_paragraph_used else text_frame.add_paragraph()
                first_paragraph_used = True
                paragraph.text = block.text or ""
                paragraph.alignment = PP_ALIGN.LEFT
                if paragraph.runs:
                    self._apply_run_style(paragraph.runs[0], self.theme.font_family, self.theme.body_font_size, self.theme.text_color)
                paragraph.space_after = Pt(14)
            elif block.type == "bullet_list":
                for item in block.items or []:
                    paragraph = text_frame.paragraphs[0] if not first_paragraph_used else text_frame.add_paragraph()
                    first_paragraph_used = True
                    paragraph.text = item
                    paragraph.level = 0
                    paragraph.alignment = PP_ALIGN.LEFT
                    if paragraph.runs:
                        self._apply_run_style(paragraph.runs[0], self.theme.font_family, self.theme.body_font_size, self.theme.text_color)
                    paragraph.space_after = Pt(8)
            elif block.type == "code":
                paragraph = text_frame.paragraphs[0] if not first_paragraph_used else text_frame.add_paragraph()
                first_paragraph_used = True
                paragraph.text = block.content or ""
                if paragraph.runs:
                    self._apply_run_style(paragraph.runs[0], "Consolas", self.theme.code_font_size, self.theme.text_color)
            elif block.type in {"table", "image", "metrics", "chart"}:
                paragraph = text_frame.paragraphs[0] if not first_paragraph_used else text_frame.add_paragraph()
                first_paragraph_used = True
                paragraph.text = f"[{block.type} content rendered elsewhere]"
                if paragraph.runs:
                    self._apply_run_style(paragraph.runs[0], self.theme.font_family, self.theme.small_font_size, self.theme.muted_text_color)

    def _add_content_panel(self, slide, left: float, top: float, width: float, height: float):
        panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, self._x(left), self._y(top), self._w(width), self._h(height))
        panel.fill.solid()
        panel.fill.fore_color.rgb = self._rgb(self.theme.card_background)
        panel.line.color.rgb = self._rgb(self.theme.border_color)
        panel.shadow.inherit = False
        panel.text_frame.margin_left = Pt(18)
        panel.text_frame.margin_right = Pt(18)
        panel.text_frame.margin_top = Pt(14)
        panel.text_frame.margin_bottom = Pt(12)
        panel.text_frame.vertical_anchor = MSO_ANCHOR.TOP
        return panel

    def _render_metric_card(self, slide, metric: MetricSpec, left_in: float, top_in: float, width_in: float, height_in: float) -> None:
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            self._x(left_in),
            self._y(top_in),
            self._w(width_in),
            self._h(height_in),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self._rgb(self.theme.card_background)
        card.line.color.rgb = self._rgb(self.theme.border_color)
        card.shadow.inherit = False

        label_box = slide.shapes.add_textbox(self._x(left_in + 0.28), self._y(top_in + 0.32), self._w(width_in - 0.5), self._h(0.4))
        paragraph = label_box.text_frame.paragraphs[0]
        paragraph.text = metric.label
        if paragraph.runs:
            self._apply_run_style(paragraph.runs[0], self.theme.font_family, 14, self.theme.muted_text_color)

        value_box = slide.shapes.add_textbox(self._x(left_in + 0.28), self._y(top_in + 1.0), self._w(width_in - 0.5), self._h(0.7))
        paragraph = value_box.text_frame.paragraphs[0]
        paragraph.text = metric.value
        if paragraph.runs:
            self._apply_run_style(paragraph.runs[0], self.theme.font_family, 26, self.theme.title_color, bold=True)

        if metric.delta:
            delta_box = slide.shapes.add_textbox(self._x(left_in + 0.28), self._y(top_in + 2.1), self._w(width_in - 0.5), self._h(0.45))
            paragraph = delta_box.text_frame.paragraphs[0]
            paragraph.text = metric.delta
            color = self.theme.accent_color if not metric.delta.startswith("-") else "B42318"
            if paragraph.runs:
                self._apply_run_style(paragraph.runs[0], self.theme.font_family, 14, color, bold=True)

    def _format_table_cell(self, cell, header: bool, stripe: bool = False) -> None:
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = self._rgb(self.theme.accent_color if header else (self.theme.accent_soft_color if stripe else self.theme.card_background))
        paragraph = cell.text_frame.paragraphs[0]
        if paragraph.runs:
            color = "FFFFFF" if header else self.theme.text_color
            self._apply_run_style(
                paragraph.runs[0],
                self.theme.font_family,
                13 if header else 12,
                color,
                bold=header,
            )

    def _apply_speaker_notes(self, slide, slide_spec: SlideSpec) -> None:
        if not slide_spec.speaker_notes:
            return
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        paragraph.text = slide_spec.speaker_notes
        if paragraph.runs:
            self._apply_run_style(paragraph.runs[0], self.theme.font_family, 12, self.theme.text_color)

    def _apply_run_style(self, run, font_name: str, font_size: int, color_hex: str, bold: bool = False) -> None:
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = self._rgb(color_hex)

    def _x(self, value: float):
        return Inches(value * self.slide_width_in / BASE_SLIDE_WIDTH)

    def _y(self, value: float):
        return Inches(value * self.slide_height_in / BASE_SLIDE_HEIGHT)

    def _w(self, value: float):
        return self._x(value)

    def _h(self, value: float):
        return self._y(value)

    @staticmethod
    def _rgb(hex_value: str) -> RGBColor:
        return RGBColor.from_string(hex_value.replace("#", "").upper())


def generate_ppt(
    spec: DeckSpec,
    output_path: str | Path,
    template_path: str | None = None,
    theme_path: str | None = None,
    template_config_path: str | None = None,
) -> ToolResult:
    renderer = PptRenderer(
        spec=spec,
        template_path=template_path,
        theme_path=theme_path,
        template_config_path=template_config_path,
    )
    return renderer.render(output_path)
