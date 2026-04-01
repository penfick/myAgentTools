from pathlib import Path
from uuid import uuid4

from pptx import Presentation

from my_agent_tools.specs import DeckSpec
from my_agent_tools.tools.ppt import generate_ppt


def _test_output(name: str) -> Path:
    output_dir = Path("out") / "test_outputs"
    output_dir.mkdir(parents=True, exist_ok=True)
    return output_dir / f"{name}_{uuid4().hex}.pptx"


def test_generate_ppt_writes_file():
    spec = DeckSpec.from_path("examples/sample_deck.json")
    output = _test_output("sample")

    result = generate_ppt(spec, output_path=output)

    assert result.ok is True
    assert output.exists()
    assert output.stat().st_size > 0


def test_generate_ppt_supports_chart_and_pagination():
    spec = DeckSpec.model_validate(
        {
            "meta": {
                "title": "Stress Deck",
                "ratio": "16:9",
            },
            "constraints": {
                "max_text_line_units_per_slide": 6,
                "max_chars_per_text_slide": 180,
                "max_table_rows_per_slide": 4,
            },
            "slides": [
                {
                    "title": "Long Summary",
                    "kind": "content",
                    "blocks": [
                        {
                            "type": "bullet_list",
                            "items": [f"Bullet item number {index} with enough detail to require pagination." for index in range(1, 9)],
                        }
                    ],
                },
                {
                    "title": "Large Table",
                    "kind": "table",
                    "blocks": [
                        {
                            "type": "table",
                            "columns": ["Region", "Revenue", "QoQ"],
                            "rows": [[f"R{index}", index * 10, f"+{index}%"] for index in range(1, 13)],
                        }
                    ],
                },
                {
                    "title": "Trend",
                    "kind": "chart",
                    "blocks": [
                        {
                            "type": "chart",
                            "chart": {
                                "chart_type": "column_clustered",
                                "categories": ["Jan", "Feb", "Mar"],
                                "series": [
                                    {"name": "2025", "values": [12, 18, 22]},
                                    {"name": "2026", "values": [16, 24, 29]},
                                ],
                            },
                        }
                    ],
                },
            ],
        }
    )
    output = _test_output("stress")

    result = generate_ppt(spec, output_path=output)

    assert result.ok is True
    presentation = Presentation(output)
    assert len(presentation.slides) > len(spec.slides)
    assert any(shape.has_chart for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "has_chart"))


def test_generate_ppt_supports_template_mapping():
    spec = DeckSpec.model_validate(
        {
            "meta": {
                "title": "Template Deck",
                "subtitle": "Built-in layouts",
            },
            "constraints": {
                "template_config_path": "examples/default_template_config.json",
            },
            "slides": [
                {
                    "title": "Template Deck",
                    "subtitle": "Built-in layouts",
                    "kind": "title",
                },
                {
                    "title": "Overview",
                    "kind": "content",
                    "blocks": [
                        {
                            "type": "bullet_list",
                            "items": ["First point", "Second point"],
                        }
                    ],
                },
            ],
        }
    )
    output = _test_output("template")

    result = generate_ppt(spec, output_path=output)

    assert result.ok is True
    assert result.meta["template_mode"] is True
    presentation = Presentation(output)
    assert presentation.slides[0].shapes.title.text == "Template Deck"
    body_placeholder = presentation.slides[1].placeholders[1]
    assert "First point" in body_placeholder.text
